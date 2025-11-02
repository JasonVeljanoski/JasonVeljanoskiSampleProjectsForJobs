import os
from io import BytesIO

from docx.shared import Mm
from docxtpl import DocxTemplate, InlineImage
from fastapi import Depends
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.parts.image import Image
from pptx.util import Inches, Pt

from app import schemas, utils
from app.models import five_why

from ..schemas.enums import DocumentPaths, EventTypeEnum, ImgPath, TemplateDocNames
from .date import format_date


# ------------------
# .docx TEMPLATES
# ------------------
def render_docxtemplate(context, template_path, out_path):
    """
    Apply context to office document and output newly created file in out_path
    """
    tpl = DocxTemplate(template_path)
    tpl.render(context)
    tpl.save(out_path)


def get_user(*, db=Depends(utils.get_db), id: int):
    from app import crud

    return crud.user.get(id)


# ------------------
# FIVE WHY .docx TEMPLATES
# ------------------
def build_fivewhy_context(investigation):
    def get_status(value):
        if value == schemas.enums.StatusEnum.OPEN:
            return "Open"
        elif value == schemas.enums.StatusEnum.ON_HOLD:
            return "On Hold"
        elif value == schemas.enums.StatusEnum.CLOSED:
            return "Closed"
        return "Open"

    # ROOT RESPONSE CONTEXT
    # resultant lists
    left_branch = []
    right_branch = []
    common_nodes = []

    # flags
    been_split = False
    leaf_reached = False
    left = True

    # DFS init
    stack = []
    stack.append(investigation.five_why.root_response)

    while len(stack) != 0:
        v = stack.pop()

        # is_leaf flag
        leaf_reached = False
        num_children = len(v.children_responses)
        if num_children == 0:
            leaf_reached = True

        node = {
            "cause": v.cause,
            "reason": v.reason,
            "filenames": v.filenames,
            "is_leaf": leaf_reached,
        }
        if not been_split:
            common_nodes.append(node)
        else:
            if not left:
                left_branch.append(node)
            else:
                right_branch.append(node)

        # branch and splits meta
        if num_children == 0:
            left = False
        elif num_children > 1:
            been_split = True

        # add children
        for child in v.children_responses:
            stack.append(child)

    # BUILD CONTEXT
    # construct actions object
    actions = []
    for action in investigation.five_why.actions:
        if not action.is_archived:
            actions.append(
                {
                    "title": action.title,
                    "description": action.description,
                    "date_due": format_date(action.date_due),
                    "date_closed": utils.format_date(action.date_closed),
                    "owners": action.owner_users,
                    "status": get_status(action.status),
                }
            )
    # get flash report actions
    for action in investigation.flash_report.actions:
        # add flash report actions only if been checked by user to include in report
        if action.id in investigation.five_why.flash_report_action_ids:
            actions.append(
                {
                    "title": action.title,
                    "description": action.description,
                    "date_due": format_date(action.date_due),
                    "date_closed": utils.format_date(action.date_closed),
                    "owners": action.owner_users,
                    "status": get_status(action.status),
                }
            )

    # --------------------------------------------

    context = {
        "event_description": investigation.five_why.event_description,
        "actions": actions,
        "participants": investigation.five_why.owner_users,
        "supervisor": investigation.five_why.supervisor,
        "responses": {
            "common_nodes": common_nodes,
            "left_branch": left_branch,
            "right_branch": right_branch,
        },
        "site": investigation.site,
        "department": investigation.department,
        "date": format_date(investigation.event_datetime),
    }

    return context


def render_save_five_why_template(context, template_path, out_path):
    """
    Apply context to office document and output newly created file in out_path
    """
    tpl = DocxTemplate(template_path)

    # CREATE `imgs` KEY FOR CONTEXT
    # why not doing it in create context function?
    #   you need the `tpl` object to do the below!
    for branch in [
        context["responses"]["left_branch"],
        context["responses"]["right_branch"],
    ]:
        height = Mm(37)

        for node in branch:
            res = []
            for file in node["filenames"]:
                res.append(InlineImage(tpl, f"{ImgPath.FIVE_WHY.value}/{file}", height=height))
            node["imgs"] = res

    for branch in [context["responses"]["common_nodes"]]:
        for node in branch:
            height = Mm(83)
            if len(node["filenames"]) > 1:
                height = Mm(37)
            res = []
            for file in node["filenames"]:
                res.append(InlineImage(tpl, f"{ImgPath.FIVE_WHY.value}/{file}", height=height))
            node["imgs"] = res

    tpl.render(context)
    tpl.save(out_path)


# ------------------
# FLASH REPORT .pptx TEMPLATE
# ------------------


def flashreport_pptx_template(context, out_filename):
    """
     Given a context, populate .pptx template and save output for flash report document.
     Context example:
     context = {
        "site": "Solomon",
        "department": "Kings Valley Operations",
        "title": "This is a title",
        "owner": "Jason Veljanoski",
        "event_date": "17/10/2022",
        "event_time": "9:00am",
        "event_description": "this is a description of the event",
        "effective_duration": "24",
        "lost_tonnes": "10000",
        "business_impact": "optional extra info",
        "immediate_action": "We asked for help right away",
        "root_causes": ["Potential cause uno", "Potential cause dos", "Potential cause tres"],
        "actions": [
            {
                "title": "action.title this is some longer text to see how the action loooks",
                "description": "action.description",
                "date_due": "01/01/2022",
                ...
            },
            {
                "title": "action.title",
                "description": "action.description",
                "date_due": "01/01/2022",
                ...
            },
            {
                "title": "action.title",
                "description": "action.description",
                "date_due": "01/01/2022",
                ...
            },
        ],
        "filenames": [
            'fileA.jpeg',
            'fileB.jpeg
        ]
    }
    """

    def render_slide_one(context):
        """TEMPLATING FOR SLIDE ONE"""
        MAIN_TEXT_IDX = 13
        SUB_TEXT_IDX = 23

        slide = prs.slides[0]

        text_frame = slide.placeholders[MAIN_TEXT_IDX].text_frame

        text_frame.clear()

        # create text on text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{context['site']}\n"
        run.font.size = Pt(36)
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # create text on text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{context['department']}\n"
        run.font.size = Pt(36)
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # create text on text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{context['title']}\n"
        run.font.size = Pt(36)
        run.font.color.theme_color = MSO_THEME_COLOR.DARK_2
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # create text on text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Failure Fact Sheet"
        run.font.size = Pt(36)
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        owner_strs = [f"{owner.name} - {owner.job_title}" for owner in context["owners"]]

        # !print only first name for now (requiremnts may change)
        owner = ""
        if owner:
            owner = owner_strs[0]

        slide.placeholders[SUB_TEXT_IDX].text = owner

    def render_slide_two(context):
        """TEMPLATING FOR SLIDE TWO"""
        EVENT_TITLE_IDX = 0
        WHAT_HAPPENED_ANS_IDX = 11
        IMMEDIATE_ACTIONS_ANS_IDX = 16
        ROOT_CAUSES_ANS_IDX = 14
        BUSINESS_IMPACT_ANS_IDX = 17

        IMAGE_IDX = 1
        IMAGE_TITLE_IDX = 2
        IMAGE_DESC_IDX = 3

        EQUIPMENT_IDX = 7
        OBJECT_IDX = 5
        DAMAGE_IDX = 6
        #  -------

        def paneltext_format(text_frame, text):
            """
            Formatting the main area in slide 2

            - What Happened?
            - Business Impact
            - Immediate Actions
            - Potential Root Cause(s)
            """
            # create text on text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]

            # set all line spacing to 0
            # a pre-condition for max character limit to look good
            p.space_before = Pt(0)
            p.space_after = Pt(0)

            run = p.add_run()

            # text and formating
            run.text = text
            # run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            run.font.size = Pt(12)

            # extra
            # text_frame.word_wrap = False
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        def small_text_format(text_frame, text, size=12, bold=False):
            """
            Formatting Image Title/Description Text
            """
            # create text on text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()

            # text and formating
            run.text = text
            run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            run.font.size = Pt(size)
            run.font.bold = bold
            # extra
            # text_frame.word_wrap = False
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        slide = prs.slides[1]

        # TITLE
        text_frame = slide.shapes[EVENT_TITLE_IDX].text_frame
        # create text on text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()

        # text and formating
        run.text = context["title"]
        run.font.size = Pt(24)
        # extra
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # PANEL TEXT
        paneltext_format(
            slide.shapes[WHAT_HAPPENED_ANS_IDX].text_frame, f"{context['event_description']}"
        )

        business_impact_str = ""
        event_duration = context["event_duration"] if context["event_duration"] else 0.0
        if context["event_type"] == EventTypeEnum.APLUS:
            total_effective_duration = (
                context["effective_duration"] if context["effective_duration"] else 0.0
            )
            total_tonnes_lost = context["lost_tonnes"] if context["lost_tonnes"] else 0.0
            # business_impact_str += f"{round(context['lost_tonnes'],2)} tonnes lost within {round(context['event_duration'],2)} hours"
            business_impact_str += f"Event lasted {event_duration} hours with an effective duration of {total_effective_duration} hours"
            if round(total_tonnes_lost) > 0:
                business_impact_str += f" and {round(total_tonnes_lost):,} tonnes lost"
        elif context["event_type"] == EventTypeEnum.REMS:
            business_impact_str += f"Event lasted {round(event_duration,1)} hours"

        if context["business_impact"]:
            business_impact_str += f"\n{context['business_impact']}"
        paneltext_format(
            slide.shapes[BUSINESS_IMPACT_ANS_IDX].text_frame,
            business_impact_str,
        )
        paneltext_format(
            slide.shapes[IMMEDIATE_ACTIONS_ANS_IDX].text_frame, context["immediate_action"]
        )
        paneltext_format(
            slide.shapes[ROOT_CAUSES_ANS_IDX].text_frame,
            "\n".join(
                [f"{i+1}. {context['root_causes'][i]}" for i in range(len(context["root_causes"]))]
            ),
        )

        # IMAGE TEXTS
        small_text_format(
            slide.shapes[IMAGE_TITLE_IDX].text_frame,
            f"{context['site']}, {context['department']}, {context['event_date']}",
            size=14,
            bold=True,
        )
        small_text_format(
            slide.shapes[IMAGE_DESC_IDX].text_frame, f"{context['function_location']}", bold=True
        )
        small_text_format(
            slide.shapes[EQUIPMENT_IDX].text_frame,
            f"{context['equipment']} ({context['object_type']})",
        )
        small_text_format(slide.shapes[OBJECT_IDX].text_frame, f"{context['object_part']}")
        small_text_format(slide.shapes[DAMAGE_IDX].text_frame, f"{context['damage_code']}")

        # RENDER IMAGE
        file = context["filenames"][0]  # display first only (if there are more)
        full_path = f"{ImgPath.FLASH_REPORT.value}/{file}"
        img = Image.from_file(full_path)
        # slide.shapes[IMAGE_IDX].Image = Image.from_file(full_path)
        slide.shapes[IMAGE_IDX].image._blob = img._blob

        old_pic = slide.shapes[IMAGE_IDX]
        left = old_pic.left
        top = old_pic.top
        width = old_pic.width

        new_shape = slide.shapes.add_picture(full_path, left, top, width)
        old_pic = old_pic._element
        new_pic = new_shape._element
        old_pic.addnext(new_pic)
        old_pic.getparent().remove(old_pic)

    def render_slide_three(context):
        """TEMPLATING FOR SLIDE THREE"""

        DIM_X = len(context["actions"]) + 1
        DIM_Y = 6

        slide = prs.slides[2]

        # ----------------------------
        # TITLE
        EVENT_TITLE_IDX = 0
        text_frame = slide.shapes[EVENT_TITLE_IDX].text_frame
        # create text on text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()

        # text and formating
        run.text = context["title"]
        run.font.size = Pt(24)
        # extra
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # ----------------------------

        x, y, cx, cy = Inches(0.5), Inches(2.5), Inches(11), Inches(4)
        shape = slide.shapes.add_table(DIM_X, DIM_Y, x, y, cx, cy)
        table = shape.table

        # HEADER
        table.cell(0, 0).text = "Action"
        table.cell(0, 1).text = "Responsible Person"
        table.cell(0, 2).text = "Status"
        table.cell(0, 3).text = "Priority"
        table.cell(0, 4).text = "Due Date"
        table.cell(0, 5).text = "Date Closed"

        # CONTENT
        for i, action in enumerate(context["actions"]):
            owner_strs = [f"{owner.name}" for owner in action["owners"]]
            table.cell(i + 1, 0).text = action["title"]
            table.cell(i + 1, 1).text = "\n".join(owner_strs)
            table.cell(i + 1, 2).text = action["status"]
            table.cell(i + 1, 3).text = action["priority"]
            table.cell(i + 1, 4).text = action["date_due"]
            table.cell(i + 1, 5).text = action["date_closed"]

        # STYLE
        for row in table.rows:
            row.height = Inches(0.5)

        table.columns[0].width = Inches(5)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(1)
        table.columns[3].width = Inches(1)
        table.columns[4].width = Inches(1.25)
        table.columns[5].width = Inches(1.25)

        # FONT SIZE
        def iter_cells(table):
            for row in table.rows:
                for cell in row.cells:
                    yield cell

        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

    # GLOBALS
    TEMPLATE_PATH = f"{DocumentPaths.TEMPLATES.value}/{TemplateDocNames.FLASH_REPORT_NO.value}"
    if context["sufficient_inventory_levels"]:
        TEMPLATE_PATH = f"{DocumentPaths.TEMPLATES.value}/{TemplateDocNames.FLASH_REPORT_YES.value}"

    OFFICE_PATH = f"{DocumentPaths.FLASH_REPORT.value}/{out_filename}"

    # OPEN TEMPLATE
    prs = Presentation(TEMPLATE_PATH)

    render_slide_one(context)
    render_slide_two(context)
    render_slide_three(context)

    # SAVE OUTPUT
    prs.save(OFFICE_PATH)


# --------------------------------------------


def shared_learning_pptx_template(context, out_filename):
    """
    Given a context, populate .pptx template and save output for flash report document.
    Context example:
    context = {}
    """

    def render_slide_one(context):
        """TEMPLATING FOR SLIDE ONE"""
        MAIN_TEXT_IDX = 13
        SUB_TEXT_IDX = 23

        slide = prs.slides[0]

        text_frame = slide.placeholders[MAIN_TEXT_IDX].text_frame

        text_frame.clear()

        # create text on text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{context['site']}\n"
        run.font.size = Pt(36)
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # create text on text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{context['department']}\n"
        run.font.size = Pt(36)
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # create text on text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{context['title']}\n"
        run.font.size = Pt(36)
        run.font.color.theme_color = MSO_THEME_COLOR.DARK_2
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # create text on text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "RCA Shared Learnings"
        run.font.size = Pt(36)
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        owner_strs = [f"{owner.name} - {owner.job_title}" for owner in context["owners"]]

        # !print only first name for now (requiremnts may change)
        owner = ""
        if owner:
            owner = owner_strs[0]

        slide.placeholders[SUB_TEXT_IDX].text = owner

    # -----------------------------------------

    def render_slide_two(context):
        """TEMPLATING FOR SLIDE TWO"""
        EVENT_TITLE_IDX = 0
        WHAT_HAPPENED_ANS_IDX = 6
        WHY_HAPPENED_ANS_IDX = 5
        WHY_HAPPENED_SUMMARY_IDX = 7
        SHARED_LEARNINGS_IDX = 8
        IMAGE_IDX = 4

        def paneltext_format(text_frame, text):
            """
            Formatting the main area in slide 2

            - What Happened?
            - Business Impact
            - Immediate Actions
            - Potential Root Cause(s)
            """
            # create text on text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()

            # text and formating
            run.text = text
            # run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            run.font.size = Pt(12)
            # extra
            # text_frame.word_wrap = False
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        slide = prs.slides[1]

        # TITLE
        text_frame = slide.shapes[EVENT_TITLE_IDX].text_frame
        # create text on text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()

        # text and formating
        run.text = context["title"]
        run.font.size = Pt(24)
        # extra
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # PANEL TEXT
        paneltext_format(
            slide.shapes[WHAT_HAPPENED_ANS_IDX].text_frame, f"{context['event_description']}"
        )
        why_summary_text = f"The event occurred due to {context['cause_code']} from a failure in {context['cause_category']}."
        paneltext_format(slide.shapes[WHY_HAPPENED_SUMMARY_IDX].text_frame, why_summary_text)
        paneltext_format(slide.shapes[WHY_HAPPENED_ANS_IDX].text_frame, context["reason"])
        paneltext_format(slide.shapes[SHARED_LEARNINGS_IDX].text_frame, context["shared_learnings"])

        # RENDER IMAGE
        file = context["flashreport_filename"]  # display first only (if there are more)
        full_path = f"{ImgPath.FLASH_REPORT.value}/{file}"
        img = Image.from_file(full_path)
        # slide.shapes[IMAGE_IDX].Image = Image.from_file(full_path)
        slide.shapes[IMAGE_IDX].image._blob = img._blob

        old_pic = slide.shapes[IMAGE_IDX]
        left = old_pic.left
        top = old_pic.top
        width = old_pic.width

        new_shape = slide.shapes.add_picture(full_path, left, top, width)
        old_pic = old_pic._element
        new_pic = new_shape._element
        old_pic.addnext(new_pic)
        old_pic.getparent().remove(old_pic)

    # -----------------------------------------

    def render_slide_three(context):
        """TEMPLATING FOR SLIDE THREE"""

        IMAGE_IDX = 2
        slide = prs.slides[2]

        filename = f"{context['cause_category']}.png"
        full_path = f"{DocumentPaths.TEMPLATES.value}/{filename}"
        img = Image.from_file(full_path)

        slide.shapes[IMAGE_IDX].image._blob = img._blob

        old_pic = slide.shapes[IMAGE_IDX]
        left = old_pic.left
        top = old_pic.top
        width = old_pic.width

        new_shape = slide.shapes.add_picture(full_path, left, top, width)
        old_pic = old_pic._element
        new_pic = new_shape._element
        old_pic.addnext(new_pic)
        old_pic.getparent().remove(old_pic)

    # -----------------------------------------

    def render_slide_four(context):
        """TEMPLATING FOR SLIDE FOUR"""

        DIM_X = len(context["actions"]) + 1
        DIM_Y = 6

        slide = prs.slides[3]

        # ----------------------------
        # TITLE
        EVENT_TITLE_IDX = 0
        text_frame = slide.shapes[EVENT_TITLE_IDX].text_frame
        # create text on text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()

        # text and formating
        run.text = context["title"]
        run.font.size = Pt(24)
        # extra
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # ----------------------------

        x, y, cx, cy = Inches(0.25), Inches(2), Inches(11), Inches(4)
        shape = slide.shapes.add_table(DIM_X, DIM_Y, x, y, cx, cy)
        table = shape.table

        # HEADER
        table.cell(0, 0).text = "Action"
        table.cell(0, 1).text = "Person Responsible"
        table.cell(0, 2).text = "Status"
        table.cell(0, 3).text = "Priority"
        table.cell(0, 4).text = "Due Date"
        table.cell(0, 5).text = "Date Closed"

        # CONTENT
        for i, action in enumerate(context["actions"]):
            owner_strs = [f"{owner.name}" for owner in action["owners"]]
            table.cell(i + 1, 0).text = action["title"]
            table.cell(i + 1, 1).text = "\n".join(owner_strs)
            table.cell(i + 1, 2).text = action["status"]
            table.cell(i + 1, 3).text = action["priority"]
            table.cell(i + 1, 4).text = action["date_due"]
            table.cell(i + 1, 5).text = action["date_closed"]

        # STYLE
        for row in table.rows:
            row.height = Inches(0.2)

        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(0.75)
        table.columns[3].width = Inches(1)
        table.columns[4].width = Inches(1)
        table.columns[5].width = Inches(1)

        # FONT SIZE
        def iter_cells(table):
            for row in table.rows:
                for cell in row.cells:
                    yield cell

        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)

        # RENDER IMAGE
        IMAGE_IDX = 2
        file = context["initiative_image_filename"]  # display first only (if there are more)

        if not file:
            IMAGE_TEXT_IDX = 2
            slide.shapes.element.remove(slide.shapes[IMAGE_IDX].element)
            slide.shapes.element.remove(slide.shapes[IMAGE_TEXT_IDX].element)
        else:
            full_path = f"{ImgPath.SHARED_LEARNING.value}/{file}"
            img = Image.from_file(full_path)

            slide.shapes[IMAGE_IDX].image._blob = img._blob

            old_pic = slide.shapes[IMAGE_IDX]
            left = old_pic.left
            top = old_pic.top
            width = old_pic.width

            new_shape = slide.shapes.add_picture(full_path, left, top, width)
            old_pic = old_pic._element
            new_pic = new_shape._element
            old_pic.addnext(new_pic)
            old_pic.getparent().remove(old_pic)

    # -----------------------------------------

    # GLOBALS
    TEMPLATE_PATH = f"{DocumentPaths.TEMPLATES.value}/{TemplateDocNames.SHARED_LEARNING.value}"
    OFFICE_PATH = f"{DocumentPaths.SHARED_LEARNING.value}/{out_filename}"

    # OPEN TEMPLATE
    prs = Presentation(TEMPLATE_PATH)

    render_slide_one(context)
    render_slide_two(context)
    render_slide_three(context)
    render_slide_four(context)

    # SAVE OUTPUT
    prs.save(OFFICE_PATH)
